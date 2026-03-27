"""
Browser Automation MCP Server
============================
An MCP server that gives Claude Code (or any MCP client) the ability to
control a real browser via Playwright. This bridges the gap between
CLI-based AI agents and interactive web automation.

Usage with Claude Code:
    Add to ~/.claude/settings.json or project .mcp.json:
    {
        "mcpServers": {
            "browser": {
                "command": "python",
                "args": ["/path/to/browser_mcp.py"],
                "env": {
                    "BROWSER_HEADLESS": "false"
                }
            }
        }
    }

Environment Variables:
    BROWSER_HEADLESS: "true" (default) or "false" for visible browser
    BROWSER_VIEWPORT_WIDTH: viewport width in pixels (default: 1280)
    BROWSER_VIEWPORT_HEIGHT: viewport height in pixels (default: 720)
    BROWSER_TIMEOUT: default timeout in ms (default: 30000)
    BROWSER_TYPE: "chromium" (default), "firefox", or "webkit"
"""

import asyncio
import base64
import json
import os
import re
from contextlib import asynccontextmanager
from enum import Enum
from typing import Any, Dict, List, Optional

from mcp.server.fastmcp import FastMCP, Context
from pydantic import BaseModel, Field, ConfigDict

# ─── Configuration ──────────────────────────────────────────────────────────

HEADLESS = os.getenv("BROWSER_HEADLESS", "true").lower() == "true"
VIEWPORT_WIDTH = int(os.getenv("BROWSER_VIEWPORT_WIDTH", "1280"))
VIEWPORT_HEIGHT = int(os.getenv("BROWSER_VIEWPORT_HEIGHT", "720"))
DEFAULT_TIMEOUT = int(os.getenv("BROWSER_TIMEOUT", "30000"))
BROWSER_TYPE = os.getenv("BROWSER_TYPE", "chromium")

# ─── Lazy Browser Manager ──────────────────────────────────────────────────

class _LazyBrowser:
    """Delays Playwright browser launch until the first tool call.

    This prevents the MCP server from timing out during startup, since
    launching Chromium can take several seconds on Windows.
    """

    def __init__(self):
        self.pw = None
        self.browser = None
        self.context = None
        self.page = None
        self._started = False

    async def ensure_started(self):
        """Launch the browser if it hasn't been launched yet."""
        if self._started:
            return
        from playwright.async_api import async_playwright

        self.pw = await async_playwright().start()
        launcher = getattr(self.pw, BROWSER_TYPE, self.pw.chromium)
        self.browser = await launcher.launch(headless=HEADLESS)
        self.context = await self.browser.new_context(
            viewport={"width": VIEWPORT_WIDTH, "height": VIEWPORT_HEIGHT},
            user_agent=(
                "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                "(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
            ),
        )
        self.context.set_default_timeout(DEFAULT_TIMEOUT)
        self.page = await self.context.new_page()
        self._started = True

    async def shutdown(self):
        """Clean up browser resources."""
        if not self._started:
            return
        try:
            await self.context.close()
        except Exception:
            pass
        try:
            await self.browser.close()
        except Exception:
            pass
        try:
            await self.pw.stop()
        except Exception:
            pass


@asynccontextmanager
async def browser_lifespan(server):
    """MCP lifespan — returns immediately; browser launches on first use."""
    lazy = _LazyBrowser()
    yield {"_lazy": lazy}
    await lazy.shutdown()


# ─── Initialize Server ─────────────────────────────────────────────────────

mcp = FastMCP("browser_mcp", lifespan=browser_lifespan)


# ─── Helpers ────────────────────────────────────────────────────────────────

async def _get_page(ctx: Context):
    """Get the active Playwright page, launching browser if needed."""
    lazy = ctx.request_context.lifespan_context["_lazy"]
    await lazy.ensure_started()
    return lazy.page


async def _get_context(ctx: Context):
    """Get the browser context, launching browser if needed."""
    lazy = ctx.request_context.lifespan_context["_lazy"]
    await lazy.ensure_started()
    return lazy.context


async def _resolve_locator(page, selector: str, index: int = 0):
    """
    Resolve a flexible selector to a Playwright locator.

    Supports:
      - CSS selectors: "div.class", "#id", "input[name='q']"
      - XPath: "//div[@class='example']"
      - Text: "text=Click me" or "text=Submit"
      - Role: "role=button[name='Submit']"
      - Playwright built-in: "button >> text=Submit"
    """
    locator = page.locator(selector)
    if await locator.count() == 0:
        # Fallback: try as text content match
        locator = page.get_by_text(selector, exact=False)
    if await locator.count() == 0:
        raise ValueError(
            f"No element found for selector: '{selector}'. "
            f"Try a CSS selector, XPath, or text= prefix."
        )
    return locator.nth(index)


def _truncate(text: str, max_len: int = 50000) -> str:
    """Truncate text to avoid overwhelming the LLM context."""
    if len(text) <= max_len:
        return text
    half = max_len // 2
    return text[:half] + f"\n\n... [truncated {len(text) - max_len} chars] ...\n\n" + text[-half:]


# ─── Input Models ───────────────────────────────────────────────────────────

class NavigateInput(BaseModel):
    """Input for navigating to a URL."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    url: str = Field(..., description="URL to navigate to (e.g., 'https://example.com')")
    wait_until: Optional[str] = Field(
        default="domcontentloaded",
        description="Wait condition: 'load', 'domcontentloaded', 'networkidle', or 'commit'"
    )


class ClickInput(BaseModel):
    """Input for clicking an element."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    selector: str = Field(..., description="CSS selector, XPath, text= selector, or visible text to click")
    index: int = Field(default=0, description="Which match to click if multiple elements found (0-based)", ge=0)
    button: Optional[str] = Field(default="left", description="Mouse button: 'left', 'right', or 'middle'")
    click_count: int = Field(default=1, description="Number of clicks (2 for double-click)", ge=1, le=3)


class TypeInput(BaseModel):
    """Input for typing text into an element."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    selector: str = Field(..., description="Selector for the input element to type into")
    text: str = Field(..., description="Text to type")
    index: int = Field(default=0, description="Which match to target if multiple elements found", ge=0)
    clear_first: bool = Field(default=True, description="Clear existing text before typing")
    press_enter: bool = Field(default=False, description="Press Enter after typing")


class FillInput(BaseModel):
    """Input for filling a form field (instant, no keystroke events)."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    selector: str = Field(..., description="Selector for the input element")
    value: str = Field(..., description="Value to fill")
    index: int = Field(default=0, description="Which match to target if multiple found", ge=0)


class SelectInput(BaseModel):
    """Input for selecting a dropdown option."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    selector: str = Field(..., description="Selector for the <select> element")
    value: Optional[str] = Field(default=None, description="Option value attribute to select")
    label: Optional[str] = Field(default=None, description="Option visible text to select")
    index: int = Field(default=0, description="Which match to target if multiple found", ge=0)


class HoverInput(BaseModel):
    """Input for hovering over an element."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    selector: str = Field(..., description="Selector for element to hover over")
    index: int = Field(default=0, description="Which match to target", ge=0)


class ScrollInput(BaseModel):
    """Input for scrolling the page."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    direction: str = Field(default="down", description="Scroll direction: 'up', 'down', 'left', 'right'")
    amount: int = Field(default=500, description="Pixels to scroll", ge=0, le=10000)
    selector: Optional[str] = Field(default=None, description="Optional: scroll within a specific element")


class WaitInput(BaseModel):
    """Input for waiting."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    selector: Optional[str] = Field(default=None, description="Wait for this selector to appear")
    state: Optional[str] = Field(
        default="visible",
        description="State to wait for: 'visible', 'hidden', 'attached', 'detached'"
    )
    timeout: int = Field(default=10000, description="Max wait time in ms", ge=0, le=60000)
    delay_ms: Optional[int] = Field(default=None, description="Simple delay in ms (if no selector given)", ge=0, le=30000)


class EvaluateInput(BaseModel):
    """Input for executing JavaScript in the browser."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    script: str = Field(..., description="JavaScript code to execute in the browser page context")


class FindInput(BaseModel):
    """Input for finding elements on the page."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    selector: Optional[str] = Field(default=None, description="CSS/XPath selector to find elements")
    text: Optional[str] = Field(default=None, description="Find elements containing this text")
    role: Optional[str] = Field(default=None, description="ARIA role to search for (e.g., 'button', 'link', 'textbox')")
    max_results: int = Field(default=20, description="Maximum number of results to return", ge=1, le=100)


class ScreenshotInput(BaseModel):
    """Input for taking screenshots."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    selector: Optional[str] = Field(default=None, description="Optional: screenshot a specific element instead of full page")
    full_page: bool = Field(default=False, description="Capture full scrollable page (ignored if selector given)")


class KeyboardInput(BaseModel):
    """Input for pressing keyboard keys."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    key: str = Field(
        ...,
        description=(
            "Key to press. Examples: 'Enter', 'Tab', 'Escape', 'Backspace', "
            "'ArrowDown', 'Control+a', 'Meta+c', 'Shift+Tab'"
        )
    )
    count: int = Field(default=1, description="Number of times to press the key", ge=1, le=50)


class GetTextInput(BaseModel):
    """Input for extracting text from the page."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    selector: Optional[str] = Field(
        default=None,
        description="Extract text from a specific element. If omitted, gets all page text."
    )
    index: int = Field(default=0, description="Which match to target if multiple found", ge=0)


class TabAction(str, Enum):
    """Tab management actions."""
    NEW = "new"
    CLOSE = "close"
    LIST = "list"
    SWITCH = "switch"


class TabInput(BaseModel):
    """Input for tab management."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    action: TabAction = Field(..., description="Tab action: 'new', 'close', 'list', or 'switch'")
    url: Optional[str] = Field(default=None, description="URL for new tab (action='new')")
    tab_index: Optional[int] = Field(default=None, description="Tab index for switch/close (0-based)", ge=0)


# ─── Tools ──────────────────────────────────────────────────────────────────

@mcp.tool(
    name="browser_navigate",
    annotations={
        "title": "Navigate to URL",
        "readOnlyHint": False,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": True,
    },
)
async def browser_navigate(params: NavigateInput, ctx: Context) -> str:
    """Navigate the browser to a URL. Returns the page title and final URL after navigation.

    Args:
        params (NavigateInput): Navigation parameters containing:
            - url (str): The URL to navigate to
            - wait_until (str): When to consider navigation complete

    Returns:
        str: JSON with page title and final URL
    """
    page = await _get_page(ctx)
    try:
        response = await page.goto(params.url, wait_until=params.wait_until)
        status = response.status if response else "unknown"
        return json.dumps({
            "status": "success",
            "url": page.url,
            "title": await page.title(),
            "http_status": status,
        })
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_click",
    annotations={
        "title": "Click Element",
        "readOnlyHint": False,
        "destructiveHint": False,
        "idempotentHint": False,
        "openWorldHint": True,
    },
)
async def browser_click(params: ClickInput, ctx: Context) -> str:
    """Click an element on the page. Supports CSS selectors, XPath, text matching, and more.

    Args:
        params (ClickInput): Click parameters containing:
            - selector (str): How to find the element
            - index (int): Which match to click (0-based)
            - button (str): Mouse button to use
            - click_count (int): Number of clicks

    Returns:
        str: JSON confirming click action or error details
    """
    page = await _get_page(ctx)
    try:
        locator = await _resolve_locator(page, params.selector, params.index)
        await locator.click(button=params.button, click_count=params.click_count)
        await page.wait_for_load_state("domcontentloaded")
        return json.dumps({
            "status": "success",
            "clicked": params.selector,
            "url": page.url,
            "title": await page.title(),
        })
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_type",
    annotations={
        "title": "Type Text",
        "readOnlyHint": False,
        "destructiveHint": False,
        "idempotentHint": False,
        "openWorldHint": True,
    },
)
async def browser_type(params: TypeInput, ctx: Context) -> str:
    """Type text into an input field with realistic keystroke simulation.

    Args:
        params (TypeInput): Typing parameters containing:
            - selector (str): Input element selector
            - text (str): Text to type
            - clear_first (bool): Clear existing text first
            - press_enter (bool): Press Enter after typing

    Returns:
        str: JSON confirming type action or error details
    """
    page = await _get_page(ctx)
    try:
        locator = await _resolve_locator(page, params.selector, params.index)
        if params.clear_first:
            await locator.clear()
        await locator.type(params.text, delay=50)
        if params.press_enter:
            await locator.press("Enter")
        return json.dumps({"status": "success", "typed": params.text, "selector": params.selector})
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_fill",
    annotations={
        "title": "Fill Form Field",
        "readOnlyHint": False,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": True,
    },
)
async def browser_fill(params: FillInput, ctx: Context) -> str:
    """Instantly fill a form field (no keystroke simulation). Faster than browser_type.

    Args:
        params (FillInput): Fill parameters containing:
            - selector (str): Input element selector
            - value (str): Value to fill

    Returns:
        str: JSON confirming fill or error details
    """
    page = await _get_page(ctx)
    try:
        locator = await _resolve_locator(page, params.selector, params.index)
        await locator.fill(params.value)
        return json.dumps({"status": "success", "filled": params.selector, "value": params.value})
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_select",
    annotations={
        "title": "Select Dropdown Option",
        "readOnlyHint": False,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": True,
    },
)
async def browser_select(params: SelectInput, ctx: Context) -> str:
    """Select an option from a <select> dropdown element.

    Args:
        params (SelectInput): Selection parameters containing:
            - selector (str): The <select> element selector
            - value (str): Option value attribute to select
            - label (str): Option visible text to select

    Returns:
        str: JSON with selected values or error details
    """
    page = await _get_page(ctx)
    try:
        locator = await _resolve_locator(page, params.selector, params.index)
        if params.value:
            result = await locator.select_option(value=params.value)
        elif params.label:
            result = await locator.select_option(label=params.label)
        else:
            return json.dumps({"status": "error", "message": "Provide either 'value' or 'label'"})
        return json.dumps({"status": "success", "selected": result})
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_hover",
    annotations={
        "title": "Hover Over Element",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": True,
    },
)
async def browser_hover(params: HoverInput, ctx: Context) -> str:
    """Hover the mouse over an element. Useful for revealing tooltips or dropdown menus.

    Args:
        params (HoverInput): Hover parameters containing:
            - selector (str): Element to hover over

    Returns:
        str: JSON confirming hover or error details
    """
    page = await _get_page(ctx)
    try:
        locator = await _resolve_locator(page, params.selector, params.index)
        await locator.hover()
        return json.dumps({"status": "success", "hovered": params.selector})
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_scroll",
    annotations={
        "title": "Scroll Page",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": False,
        "openWorldHint": False,
    },
)
async def browser_scroll(params: ScrollInput, ctx: Context) -> str:
    """Scroll the page or a specific element.

    Args:
        params (ScrollInput): Scroll parameters containing:
            - direction (str): 'up', 'down', 'left', 'right'
            - amount (int): Pixels to scroll
            - selector (str): Optional element to scroll within

    Returns:
        str: JSON confirming scroll action
    """
    page = await _get_page(ctx)
    dx, dy = 0, 0
    if params.direction == "down":
        dy = params.amount
    elif params.direction == "up":
        dy = -params.amount
    elif params.direction == "right":
        dx = params.amount
    elif params.direction == "left":
        dx = -params.amount

    try:
        if params.selector:
            locator = await _resolve_locator(page, params.selector)
            await locator.evaluate(f"el => el.scrollBy({dx}, {dy})")
        else:
            await page.evaluate(f"window.scrollBy({dx}, {dy})")
        return json.dumps({"status": "success", "scrolled": params.direction, "amount": params.amount})
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_wait",
    annotations={
        "title": "Wait for Element or Delay",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": False,
    },
)
async def browser_wait(params: WaitInput, ctx: Context) -> str:
    """Wait for an element to appear or a fixed time delay.

    Args:
        params (WaitInput): Wait parameters containing:
            - selector (str): Wait for this selector
            - state (str): Element state to wait for
            - timeout (int): Max wait time in ms
            - delay_ms (int): Simple delay (if no selector)

    Returns:
        str: JSON confirming wait completed or timeout error
    """
    page = await _get_page(ctx)
    try:
        if params.selector:
            await page.wait_for_selector(
                params.selector, state=params.state, timeout=params.timeout
            )
            return json.dumps({"status": "success", "found": params.selector, "state": params.state})
        elif params.delay_ms:
            await asyncio.sleep(params.delay_ms / 1000.0)
            return json.dumps({"status": "success", "waited_ms": params.delay_ms})
        else:
            return json.dumps({"status": "error", "message": "Provide either 'selector' or 'delay_ms'"})
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_screenshot",
    annotations={
        "title": "Take Screenshot",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": False,
    },
)
async def browser_screenshot(params: ScreenshotInput, ctx: Context) -> str:
    """Take a screenshot of the page or a specific element. Returns base64-encoded PNG.

    Args:
        params (ScreenshotInput): Screenshot parameters containing:
            - selector (str): Optional element to screenshot
            - full_page (bool): Capture full scrollable page

    Returns:
        str: JSON with base64 screenshot data and metadata
    """
    page = await _get_page(ctx)
    try:
        if params.selector:
            locator = await _resolve_locator(page, params.selector)
            screenshot_bytes = await locator.screenshot(type="png")
        else:
            screenshot_bytes = await page.screenshot(type="png", full_page=params.full_page)

        b64 = base64.b64encode(screenshot_bytes).decode("utf-8")
        return json.dumps({
            "status": "success",
            "format": "png",
            "encoding": "base64",
            "size_bytes": len(screenshot_bytes),
            "data": b64,
            "note": "Base64-encoded PNG image. Decode with base64.b64decode() to get raw bytes.",
        })
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_find",
    annotations={
        "title": "Find Elements",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": False,
    },
)
async def browser_find(params: FindInput, ctx: Context) -> str:
    """Find elements on the page by selector, text content, or ARIA role.
    Returns element details including tag, text, attributes, and visibility.

    Args:
        params (FindInput): Search parameters containing:
            - selector (str): CSS/XPath selector
            - text (str): Text content to search for
            - role (str): ARIA role to search for
            - max_results (int): Maximum results to return

    Returns:
        str: JSON array of matching elements with tag, text, attributes, and bounding box
    """
    page = await _get_page(ctx)
    try:
        if params.selector:
            locator = page.locator(params.selector)
        elif params.text:
            locator = page.get_by_text(params.text, exact=False)
        elif params.role:
            locator = page.get_by_role(params.role)
        else:
            return json.dumps({"status": "error", "message": "Provide selector, text, or role"})

        count = await locator.count()
        count = min(count, params.max_results)

        elements = []
        for i in range(count):
            el = locator.nth(i)
            try:
                tag = await el.evaluate("el => el.tagName.toLowerCase()")
                text = (await el.inner_text())[:200] if await el.is_visible() else ""
                attrs = await el.evaluate("""el => {
                    const result = {};
                    for (const attr of el.attributes) {
                        if (['id', 'class', 'name', 'type', 'href', 'src', 'placeholder',
                             'aria-label', 'role', 'value', 'alt', 'title'].includes(attr.name)) {
                            result[attr.name] = attr.value;
                        }
                    }
                    return result;
                }""")
                visible = await el.is_visible()
                bbox = await el.bounding_box() if visible else None
                elements.append({
                    "index": i,
                    "tag": tag,
                    "text": text.strip(),
                    "attributes": attrs,
                    "visible": visible,
                    "bounding_box": bbox,
                })
            except Exception:
                continue

        return json.dumps({
            "status": "success",
            "total_found": await locator.count(),
            "returned": len(elements),
            "elements": elements,
        }, indent=2)
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_get_text",
    annotations={
        "title": "Extract Page Text",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": False,
    },
)
async def browser_get_text(params: GetTextInput, ctx: Context) -> str:
    """Extract text content from the page or a specific element.

    Args:
        params (GetTextInput): Parameters containing:
            - selector (str): Optional element to extract text from
            - index (int): Which match to target

    Returns:
        str: JSON with extracted text content, page URL, and title
    """
    page = await _get_page(ctx)
    try:
        if params.selector:
            locator = await _resolve_locator(page, params.selector, params.index)
            text = await locator.inner_text()
        else:
            text = await page.inner_text("body")

        return json.dumps({
            "status": "success",
            "url": page.url,
            "title": await page.title(),
            "text": _truncate(text),
        })
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_evaluate",
    annotations={
        "title": "Execute JavaScript",
        "readOnlyHint": False,
        "destructiveHint": False,
        "idempotentHint": False,
        "openWorldHint": True,
    },
)
async def browser_evaluate(params: EvaluateInput, ctx: Context) -> str:
    """Execute JavaScript code in the browser page context. Returns the result.

    Args:
        params (EvaluateInput): Parameters containing:
            - script (str): JavaScript code to execute

    Returns:
        str: JSON with the JavaScript return value
    """
    page = await _get_page(ctx)
    try:
        result = await page.evaluate(params.script)
        return json.dumps({"status": "success", "result": result}, default=str)
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_keyboard",
    annotations={
        "title": "Press Keyboard Key",
        "readOnlyHint": False,
        "destructiveHint": False,
        "idempotentHint": False,
        "openWorldHint": False,
    },
)
async def browser_keyboard(params: KeyboardInput, ctx: Context) -> str:
    """Press a keyboard key or shortcut. Supports modifiers like Control+a, Meta+c.

    Args:
        params (KeyboardInput): Parameters containing:
            - key (str): Key to press (e.g., 'Enter', 'Tab', 'Control+a')
            - count (int): Number of times to press

    Returns:
        str: JSON confirming key press
    """
    page = await _get_page(ctx)
    try:
        for _ in range(params.count):
            await page.keyboard.press(params.key)
        return json.dumps({"status": "success", "key": params.key, "count": params.count})
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_back",
    annotations={
        "title": "Go Back",
        "readOnlyHint": False,
        "destructiveHint": False,
        "idempotentHint": False,
        "openWorldHint": True,
    },
)
async def browser_back(ctx: Context) -> str:
    """Navigate back in browser history.

    Returns:
        str: JSON with new page URL and title
    """
    page = await _get_page(ctx)
    try:
        await page.go_back(wait_until="domcontentloaded")
        return json.dumps({"status": "success", "url": page.url, "title": await page.title()})
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_forward",
    annotations={
        "title": "Go Forward",
        "readOnlyHint": False,
        "destructiveHint": False,
        "idempotentHint": False,
        "openWorldHint": True,
    },
)
async def browser_forward(ctx: Context) -> str:
    """Navigate forward in browser history.

    Returns:
        str: JSON with new page URL and title
    """
    page = await _get_page(ctx)
    try:
        await page.go_forward(wait_until="domcontentloaded")
        return json.dumps({"status": "success", "url": page.url, "title": await page.title()})
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_tabs",
    annotations={
        "title": "Manage Tabs",
        "readOnlyHint": False,
        "destructiveHint": False,
        "idempotentHint": False,
        "openWorldHint": True,
    },
)
async def browser_tabs(params: TabInput, ctx: Context) -> str:
    """Manage browser tabs: create new, close, list, or switch between tabs.

    Args:
        params (TabInput): Tab management parameters containing:
            - action (str): 'new', 'close', 'list', or 'switch'
            - url (str): URL for new tab
            - tab_index (int): Tab index for switch/close

    Returns:
        str: JSON with tab operation results
    """
    browser_ctx = await _get_context(ctx)
    lazy = ctx.request_context.lifespan_context["_lazy"]

    try:
        pages = browser_ctx.pages

        if params.action == TabAction.LIST:
            tabs = []
            for i, p in enumerate(pages):
                tabs.append({
                    "index": i,
                    "url": p.url,
                    "title": await p.title(),
                    "is_active": p == lazy.page,
                })
            return json.dumps({"status": "success", "tabs": tabs})

        elif params.action == TabAction.NEW:
            new_page = await browser_ctx.new_page()
            if params.url:
                await new_page.goto(params.url, wait_until="domcontentloaded")
            lazy.page = new_page
            return json.dumps({
                "status": "success",
                "action": "new_tab",
                "tab_index": len(browser_ctx.pages) - 1,
                "url": new_page.url,
            })

        elif params.action == TabAction.SWITCH:
            if params.tab_index is None or params.tab_index >= len(pages):
                return json.dumps({
                    "status": "error",
                    "message": f"Invalid tab index. Available: 0-{len(pages)-1}",
                })
            lazy.page = pages[params.tab_index]
            await pages[params.tab_index].bring_to_front()
            return json.dumps({
                "status": "success",
                "action": "switched",
                "tab_index": params.tab_index,
                "url": pages[params.tab_index].url,
            })

        elif params.action == TabAction.CLOSE:
            if len(pages) <= 1:
                return json.dumps({"status": "error", "message": "Cannot close the last tab"})
            idx = params.tab_index if params.tab_index is not None else len(pages) - 1
            if idx >= len(pages):
                return json.dumps({"status": "error", "message": f"Invalid tab index: {idx}"})
            closing = pages[idx]
            was_active = closing == lazy.page
            await closing.close()
            if was_active:
                lazy.page = browser_ctx.pages[min(idx, len(browser_ctx.pages) - 1)]
            return json.dumps({"status": "success", "action": "closed", "closed_index": idx})

    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_page_info",
    annotations={
        "title": "Get Page Info",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": False,
    },
)
async def browser_page_info(ctx: Context) -> str:
    """Get current page information: URL, title, viewport size, and basic page stats.

    Returns:
        str: JSON with page URL, title, viewport dimensions, and element counts
    """
    page = await _get_page(ctx)
    try:
        info = await page.evaluate("""() => ({
            url: window.location.href,
            title: document.title,
            viewport: { width: window.innerWidth, height: window.innerHeight },
            scroll: { x: window.scrollX, y: window.scrollY,
                      maxX: document.body.scrollWidth - window.innerWidth,
                      maxY: document.body.scrollHeight - window.innerHeight },
            counts: {
                links: document.querySelectorAll('a').length,
                buttons: document.querySelectorAll('button').length,
                inputs: document.querySelectorAll('input, textarea, select').length,
                images: document.querySelectorAll('img').length,
                forms: document.querySelectorAll('form').length,
            }
        })""")
        return json.dumps({"status": "success", **info})
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="browser_get_html",
    annotations={
        "title": "Get Page HTML",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": False,
    },
)
async def browser_get_html(ctx: Context) -> str:
    """Get the current page's HTML content (truncated for context management).

    Returns:
        str: The page HTML content
    """
    page = await _get_page(ctx)
    try:
        html = await page.content()
        return _truncate(html, 80000)
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


# ─── File Processing Input Models ────────────────────────────────────────


class FilePathInput(BaseModel):
    """Input requiring a file path."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    path: str = Field(..., description="Absolute path to the file")


class ExcelReadInput(BaseModel):
    """Input for reading an Excel file."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    path: str = Field(..., description="Absolute path to the .xlsx or .xls file")
    sheet: Optional[str] = Field(default=None, description="Sheet name to read. If omitted, reads the active/first sheet.")
    max_rows: int = Field(default=200, description="Maximum number of data rows to return", ge=1, le=10000)
    start_row: int = Field(default=1, description="Row number to start reading from (1-based)", ge=1)


class CsvReadInput(BaseModel):
    """Input for reading a CSV file."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    path: str = Field(..., description="Absolute path to the .csv file")
    max_rows: int = Field(default=200, description="Maximum number of data rows to return", ge=1, le=10000)
    delimiter: str = Field(default=",", description="Column delimiter character")
    encoding: str = Field(default="utf-8", description="File encoding (e.g. utf-8, latin-1, cp1252)")


class ImageReadInput(BaseModel):
    """Input for reading an image file."""
    model_config = ConfigDict(str_strip_whitespace=True, extra="forbid")
    path: str = Field(..., description="Absolute path to the image file (JPG, PNG, GIF, BMP, WebP, TIFF)")
    max_dimension: int = Field(default=1568, description="Resize the longest side to this many pixels if larger (saves tokens)", ge=100, le=4000)


# ─── File Processing Tools ───────────────────────────────────────────────


@mcp.tool(
    name="file_info",
    annotations={
        "title": "Get File Info",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": False,
    },
)
async def file_info(params: FilePathInput, ctx: Context) -> str:
    """Get file metadata: size, type, modification date, and whether it's readable.

    Args:
        params (FilePathInput): Contains the file path.

    Returns:
        str: JSON with file size, extension, modified time, and readability.
    """
    from datetime import datetime as dt

    try:
        p = os.path.abspath(params.path)
        if not os.path.exists(p):
            return json.dumps({"status": "error", "message": f"File not found: {p}"})

        st = os.stat(p)
        ext = os.path.splitext(p)[1].lower()

        supported_types = {
            ".xlsx": "Excel Workbook", ".xls": "Excel Workbook (legacy)",
            ".csv": "CSV", ".tsv": "TSV",
            ".docx": "Word Document", ".doc": "Word Document (legacy)",
            ".pptx": "PowerPoint Presentation", ".ppt": "PowerPoint (legacy)",
            ".jpg": "JPEG Image", ".jpeg": "JPEG Image", ".png": "PNG Image",
            ".gif": "GIF Image", ".bmp": "BMP Image", ".webp": "WebP Image",
            ".tiff": "TIFF Image", ".tif": "TIFF Image",
            ".pdf": "PDF Document",
            ".txt": "Text File", ".md": "Markdown File", ".json": "JSON File",
        }

        return json.dumps({
            "status": "success",
            "path": p,
            "name": os.path.basename(p),
            "extension": ext,
            "file_type": supported_types.get(ext, "Unknown"),
            "size_bytes": st.st_size,
            "size_human": f"{st.st_size / 1024:.1f} KB" if st.st_size < 1048576 else f"{st.st_size / 1048576:.1f} MB",
            "modified": dt.fromtimestamp(st.st_mtime).isoformat(),
            "readable": os.access(p, os.R_OK),
        })
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="file_list_sheets",
    annotations={
        "title": "List Excel Sheets",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": False,
    },
)
async def file_list_sheets(params: FilePathInput, ctx: Context) -> str:
    """List all sheet names in an Excel workbook.

    Args:
        params (FilePathInput): Contains the path to the Excel file.

    Returns:
        str: JSON with list of sheet names and row/column counts per sheet.
    """
    try:
        from openpyxl import load_workbook
        wb = load_workbook(params.path, read_only=True, data_only=True)
        sheets = []
        for name in wb.sheetnames:
            ws = wb[name]
            sheets.append({
                "name": name,
                "rows": ws.max_row,
                "columns": ws.max_column,
            })
        wb.close()
        return json.dumps({"status": "success", "path": params.path, "sheets": sheets})
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="file_read_excel",
    annotations={
        "title": "Read Excel File",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": False,
    },
)
async def file_read_excel(params: ExcelReadInput, ctx: Context) -> str:
    """Read an Excel file and return contents as a markdown table.

    Args:
        params (ExcelReadInput): Contains path, optional sheet name, max_rows, and start_row.

    Returns:
        str: JSON with sheet name, row/column counts, headers, and data as a markdown table.
    """
    try:
        from openpyxl import load_workbook
        wb = load_workbook(params.path, read_only=True, data_only=True)

        if params.sheet:
            if params.sheet not in wb.sheetnames:
                wb.close()
                return json.dumps({"status": "error", "message": f"Sheet '{params.sheet}' not found. Available: {wb.sheetnames}"})
            ws = wb[params.sheet]
        else:
            ws = wb.active

        sheet_title = ws.title
        rows = list(ws.iter_rows(values_only=True))
        wb.close()

        if not rows:
            return json.dumps({"status": "success", "sheet": sheet_title, "message": "Sheet is empty", "table": ""})

        headers = [str(h) if h is not None else "" for h in rows[0]]

        start_idx = params.start_row
        end_idx = start_idx + params.max_rows
        data_rows = rows[start_idx:end_idx]
        total_data_rows = len(rows) - 1

        md_lines = []
        md_lines.append("| " + " | ".join(headers) + " |")
        md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for row in data_rows:
            cells = [str(c) if c is not None else "" for c in row]
            while len(cells) < len(headers):
                cells.append("")
            cells = cells[:len(headers)]
            md_lines.append("| " + " | ".join(cells) + " |")

        table = "\n".join(md_lines)

        return json.dumps({
            "status": "success",
            "sheet": sheet_title,
            "total_rows": total_data_rows,
            "rows_returned": len(data_rows),
            "columns": len(headers),
            "headers": headers,
            "truncated": total_data_rows > len(data_rows),
            "table": _truncate(table, 80000),
        })
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="file_read_csv",
    annotations={
        "title": "Read CSV File",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": False,
    },
)
async def file_read_csv(params: CsvReadInput, ctx: Context) -> str:
    """Read a CSV file and return contents as a markdown table.

    Args:
        params (CsvReadInput): Contains path, max_rows, delimiter, and encoding.

    Returns:
        str: JSON with headers, row counts, and data as a markdown table.
    """
    import csv

    try:
        with open(params.path, "r", encoding=params.encoding, newline="") as f:
            reader = csv.reader(f, delimiter=params.delimiter)
            rows = []
            for i, row in enumerate(reader):
                if i > params.max_rows:
                    break
                rows.append(row)

        if not rows:
            return json.dumps({"status": "success", "message": "File is empty", "table": ""})

        headers = rows[0]
        data_rows = rows[1:]

        md_lines = []
        md_lines.append("| " + " | ".join(headers) + " |")
        md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for row in data_rows:
            while len(row) < len(headers):
                row.append("")
            row = row[:len(headers)]
            md_lines.append("| " + " | ".join(row) + " |")

        table = "\n".join(md_lines)

        return json.dumps({
            "status": "success",
            "total_rows": len(data_rows),
            "columns": len(headers),
            "headers": headers,
            "truncated": len(data_rows) >= params.max_rows,
            "table": _truncate(table, 80000),
        })
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="file_read_word",
    annotations={
        "title": "Read Word Document",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": False,
    },
)
async def file_read_word(params: FilePathInput, ctx: Context) -> str:
    """Read a .docx Word document and return its text content, including tables.

    Args:
        params (FilePathInput): Contains the path to the .docx file.

    Returns:
        str: JSON with document text, table count, and paragraph count.
    """
    try:
        from docx import Document
        doc = Document(params.path)

        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                if para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "").replace("Heading", "1")
                    try:
                        level = int(level)
                    except ValueError:
                        level = 1
                    paragraphs.append("#" * level + " " + text)
                else:
                    paragraphs.append(text)

        tables_md = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                rows.insert(1, header_sep)
                tables_md.append("\n".join(rows))

        content = "\n\n".join(paragraphs)
        if tables_md:
            content += "\n\n---\n\n## Tables\n\n" + "\n\n".join(tables_md)

        return json.dumps({
            "status": "success",
            "path": params.path,
            "paragraphs": len(paragraphs),
            "tables": len(doc.tables),
            "content": _truncate(content, 80000),
        })
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="file_read_powerpoint",
    annotations={
        "title": "Read PowerPoint Presentation",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": False,
    },
)
async def file_read_powerpoint(params: FilePathInput, ctx: Context) -> str:
    """Read a .pptx PowerPoint file and return text from all slides.

    Args:
        params (FilePathInput): Contains the path to the .pptx file.

    Returns:
        str: JSON with slide count and text content organized by slide.
    """
    try:
        from pptx import Presentation
        prs = Presentation(params.path)

        slides_content = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append("| " + " | ".join(cells) + " |")
                    if rows:
                        header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                        rows.insert(1, header_sep)
                        texts.append("\n".join(rows))

            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            slide_text = "\n".join(texts)
            slide_entry = f"## Slide {i}\n\n{slide_text}"
            if notes:
                slide_entry += f"\n\n**Notes:** {notes}"
            slides_content.append(slide_entry)

        content = "\n\n---\n\n".join(slides_content)

        return json.dumps({
            "status": "success",
            "path": params.path,
            "total_slides": len(prs.slides),
            "content": _truncate(content, 80000),
        })
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


@mcp.tool(
    name="file_read_image",
    annotations={
        "title": "Read Image File",
        "readOnlyHint": True,
        "destructiveHint": False,
        "idempotentHint": True,
        "openWorldHint": False,
    },
)
async def file_read_image(params: ImageReadInput, ctx: Context) -> str:
    """Read an image file and return it as base64 for Claude's vision, plus metadata.

    Supports: JPG, PNG, GIF, BMP, WebP, TIFF. Large images are automatically
    resized to save context tokens.

    Args:
        params (ImageReadInput): Contains file path and optional max_dimension for resizing.

    Returns:
        str: JSON with base64-encoded image data, dimensions, format, and file size.
    """
    try:
        from PIL import Image
        import io

        p = os.path.abspath(params.path)
        if not os.path.exists(p):
            return json.dumps({"status": "error", "message": f"File not found: {p}"})

        img = Image.open(p)
        original_size = img.size
        img_format = img.format or "PNG"

        if img.mode in ("RGBA", "P", "LA"):
            img = img.convert("RGBA")
            background = Image.new("RGBA", img.size, (255, 255, 255, 255))
            background.paste(img, mask=img.split()[3])
            img = background.convert("RGB")
        elif img.mode not in ("RGB", "L"):
            img = img.convert("RGB")

        max_dim = params.max_dimension
        w, h = img.size
        if max(w, h) > max_dim:
            ratio = max_dim / max(w, h)
            img = img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)

        buf = io.BytesIO()
        img.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode("utf-8")

        return json.dumps({
            "status": "success",
            "path": p,
            "original_size": {"width": original_size[0], "height": original_size[1]},
            "returned_size": {"width": img.size[0], "height": img.size[1]},
            "original_format": img_format,
            "encoding": "base64",
            "media_type": "image/png",
            "size_bytes": len(buf.getvalue()),
            "data": b64,
        })
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})


# ─── Entry Point ────────────────────────────────────────────────────────────

if __name__ == "__main__":
    mcp.run()
